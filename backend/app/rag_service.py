import math
import json
from typing import List, Tuple
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from . import ai_service, models, crud
from sqlalchemy.orm import Session

def extract_text_from_pdf(filepath: str) -> str:
    try:
        reader = PdfReader(filepath)
        text = ""
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
        return text
    except Exception as e:
        print(f"Error reading PDF: {e}")
        return ""

def extract_text_from_docx(filepath: str) -> str:
    try:
        doc = DocxDocument(filepath)
        text = []
        for paragraph in doc.paragraphs:
            if paragraph.text:
                text.append(paragraph.text)
        return "\n".join(text)
    except Exception as e:
        print(f"Error reading DOCX: {e}")
        return ""

def extract_text_from_pptx(filepath: str) -> str:
    try:
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text.append(run.text)
        return "\n".join(text)
    except Exception as e:
        print(f"Error reading PPTX: {e}")
        return ""

def extract_text_from_file(filepath: str, filename: str) -> str:
    ext = filename.split(".")[-1].lower()
    if ext == "pdf":
        return extract_text_from_pdf(filepath)
    elif ext == "docx":
        return extract_text_from_docx(filepath)
    elif ext in ["pptx", "ppt"]:
        return extract_text_from_pptx(filepath)
    else:
        # Fallback to UTF-8 text file reading
        try:
            with open(filepath, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            print(f"Error reading text file: {e}")
            return ""

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    """
    Split text into overlapping chunks.
    """
    if not text:
        return []
    words = text.split()
    chunks = []
    
    # Simple sliding window over words
    # A chunk_size of 1000 characters is roughly 150-200 words.
    # Let's chunk by character count for precision.
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end])
        start += (chunk_size - overlap)
        
    return chunks

def ingest_document(db: Session, filepath: str, filename: str, title: str) -> models.Document:
    """
    Extracts content, extracts concepts via Gemini, chunks text, generates embeddings, and saves to database.
    """
    content = extract_text_from_file(filepath, filename)
    if not content.strip():
        raise ValueError("Failed to extract readable text from document.")
        
    # Extract concepts
    try:
        concepts_json = ai_service.extract_concepts(content)
        # Verify it parses as JSON
        json.loads(concepts_json)
    except Exception as e:
        print(f"Error extracting concepts, using fallback empty list: {e}")
        concepts_json = "[]"
        
    # Create Document record
    doc_schema = schemas = type('schemas', (), {})()
    doc_schema.title = title
    doc_schema.filename = filename
    doc_schema.content = content
    doc_schema.concepts = concepts_json
    
    db_doc = crud.create_document(db, doc_schema)
    
    # Chunk and embed
    chunks = chunk_text(content)
    chunk_embeddings = []
    for chunk in chunks:
        # Generate embedding via Gemini
        emb = ai_service.generate_embeddings(chunk)
        chunk_embeddings.append((chunk, emb))
        
    crud.create_document_chunks(db, db_doc.id, chunk_embeddings)
    
    # Save concepts in the Spaced Repetition Tracker automatically
    try:
        concepts_list = json.loads(concepts_json)
        for c in concepts_list:
            if isinstance(c, dict) and "concept" in c:
                item_schema = type('item_schema', (), {})()
                item_schema.concept_name = c["concept"]
                item_schema.subject = title
                crud.create_spaced_repetition_item(db, item_schema)
    except Exception as e:
        print(f"Failed to auto-schedule spaced repetition: {e}")
        
    return db_doc

def cosine_similarity(v1: List[float], v2: List[float]) -> float:
    """
    Calculates cosine similarity between two vectors. Pure Python to avoid compiling native libraries.
    """
    if len(v1) != len(v2):
        return 0.0
    dot_product = sum(a * b for a, b in zip(v1, v2))
    sum1 = sum(a * a for a in v1)
    sum2 = sum(b * b for b in v2)
    if not sum1 or not sum2:
        return 0.0
    return dot_product / (math.sqrt(sum1) * math.sqrt(sum2))

def search_relevant_chunks(db: Session, query: str, document_id: int = None, top_k: int = 4) -> List[Tuple[models.DocumentChunk, float]]:
    """
    Search chunks similar to query using local cosine similarity on Gemini embeddings.
    """
    query_emb = ai_service.generate_query_embedding(query)
    
    if document_id:
        chunks = crud.get_document_chunks(db, document_id)
    else:
        chunks = crud.get_all_document_chunks(db)
        
    results = []
    for chunk in chunks:
        try:
            chunk_emb = json.loads(chunk.embedding)
            sim = cosine_similarity(query_emb, chunk_emb)
            results.append((chunk, sim))
        except Exception as e:
            print(f"Error parsing embedding for chunk {chunk.id}: {e}")
            
    # Sort by similarity descending
    results.sort(key=lambda x: x[1], reverse=True)
    return results[:top_k]
